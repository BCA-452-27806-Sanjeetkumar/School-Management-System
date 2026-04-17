#!/usr/bin/env python3
"""
School Hub Dashboard - Professional PPT Generator
Creates a comprehensive presentation with all required sections
Author: AI Assistant
Date: February 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

# Color Scheme - Professional Blue & Teal
PRIMARY_COLOR = RGBColor(13, 71, 161)      # Deep Blue
SECONDARY_COLOR = RGBColor(0, 150, 136)   # Teal
ACCENT_COLOR = RGBColor(255, 152, 0)      # Orange
TEXT_COLOR = RGBColor(33, 33, 33)         # Dark Gray
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 248, 255)        # Alice Blue

class SchoolDashboardPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.slide_number = 0
        
    def add_title_slide(self, title, subtitle, date_text="February 13-14, 2026"):
        """Add a professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        self.slide_number += 1
        
        # Add gradient-like background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Add decorative shape at top
        top_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(2.5)
        )
        top_shape.fill.solid()
        top_shape.fill.fore_color.rgb = SECONDARY_COLOR
        top_shape.line.color.rgb = SECONDARY_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1),
            Inches(9), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(66)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle with light background
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(2)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = PRIMARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(2), Inches(6),
            Inches(6), Inches(1)
        )
        date_frame = date_box.text_frame
        date_frame.text = date_text
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(20)
        date_para.font.color.rgb = SECONDARY_COLOR
        date_para.alignment = PP_ALIGN.CENTER
        
    def add_content_slide(self, title, content_list):
        """Add a slide with title and bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header shape
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Content area with light background
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = LIGHT_BG
        content_shape.line.color.rgb = SECONDARY_COLOR
        content_shape.line.width = Pt(2)
        
        # Add content
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(8), Inches(5.2)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        para_idx = 0
        for item in content_list:
            if isinstance(item, tuple):
                text, level = item
            else:
                text = str(item)
                level = 0
            
            # Skip if text is just a number
            if isinstance(text, int):
                continue
            
            if para_idx > 0:
                text_frame.add_paragraph()
            
            p = text_frame.paragraphs[para_idx]
            p.text = str(text)
            p.font.size = Pt(18 - level * 2)
            p.font.color.rgb = TEXT_COLOR
            p.level = level
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            
            # Add bullet point styling
            if level == 0:
                p.font.bold = True
                p.font.color.rgb = SECONDARY_COLOR
            
            para_idx += 1
        
        return slide
    
    def add_two_column_slide(self, title, left_content, right_content):
        """Add a slide with two columns"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Left column
        left_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(4.5), Inches(5.8)
        )
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = LIGHT_BG
        left_shape.line.color.rgb = SECONDARY_COLOR
        left_shape.line.width = Pt(2)
        
        left_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(3.9), Inches(5.2)
        )
        left_frame = left_text.text_frame
        left_frame.word_wrap = True
        for i, text in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        # Right column
        right_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(1.2),
            Inches(4.5), Inches(5.8)
        )
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = LIGHT_BG
        right_shape.line.color.rgb = SECONDARY_COLOR
        right_shape.line.width = Pt(2)
        
        right_text = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.5),
            Inches(3.9), Inches(5.2)
        )
        right_frame = right_text.text_frame
        right_frame.word_wrap = True
        for i, text in enumerate(right_content):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    def add_gantt_chart_slide(self):
        """Add Gantt Chart slide for 3-month timeline"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Project Timeline - 3 Month Duration"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Timeline content
        timeline_data = [
            ("Month 1: Planning & Design", "Weeks 1-4", ["Requirements Analysis", "Database Design", "Architecture Planning"]),
            ("Month 2: Development", "Weeks 5-8", ["Frontend Development", "Backend API", "Integration Testing"]),
            ("Month 3: Testing & Deployment", "Weeks 9-12", ["QA Testing", "Bug Fixes", "Production Deployment"]),
        ]
        
        y_start = 1.5
        for idx, (title, weeks, tasks) in enumerate(timeline_data):
            # Phase box
            phase_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y_start + idx * 1.8),
                Inches(9), Inches(1.5)
            )
            
            colors = [SECONDARY_COLOR, ACCENT_COLOR, RGBColor(76, 175, 80)]
            phase_box.fill.solid()
            phase_box.fill.fore_color.rgb = colors[idx]
            phase_box.line.color.rgb = colors[idx]
            
            # Phase title
            phase_title = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_start + idx * 1.8 + 0.1),
                Inches(3), Inches(0.3)
            )
            phase_frame = phase_title.text_frame
            phase_frame.text = title
            phase_para = phase_frame.paragraphs[0]
            phase_para.font.size = Pt(16)
            phase_para.font.bold = True
            phase_para.font.color.rgb = WHITE
            
            # Weeks
            weeks_box = slide.shapes.add_textbox(
                Inches(4), Inches(y_start + idx * 1.8 + 0.1),
                Inches(2), Inches(0.3)
            )
            weeks_frame = weeks_box.text_frame
            weeks_frame.text = weeks
            weeks_para = weeks_frame.paragraphs[0]
            weeks_para.font.size = Pt(14)
            weeks_para.font.color.rgb = WHITE
            
            # Tasks
            tasks_text = ", ".join(tasks)
            tasks_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_start + idx * 1.8 + 0.5),
                Inches(8.4), Inches(0.8)
            )
            tasks_frame = tasks_box.text_frame
            tasks_frame.text = tasks_text
            tasks_frame.word_wrap = True
            tasks_para = tasks_frame.paragraphs[0]
            tasks_para.font.size = Pt(12)
            tasks_para.font.color.rgb = WHITE
        
        # Key dates box
        dates_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.2),
            Inches(9), Inches(0.9)
        )
        dates_box.fill.solid()
        dates_box.fill.fore_color.rgb = LIGHT_BG
        dates_box.line.color.rgb = PRIMARY_COLOR
        dates_box.line.width = Pt(2)
        
        dates_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.3),
            Inches(8.4), Inches(0.7)
        )
        dates_frame = dates_text.text_frame
        dates_frame.word_wrap = True
        dates_frame.text = "📅 Level-1: 13-14 Feb 2026 | 📅 Level-2: 27-28 Feb 2026 | 📅 Level-3: 11-14 Mar 2026 | 📅 Level-4: 25-28 Mar 2026"
        dates_para = dates_frame.paragraphs[0]
        dates_para.font.size = Pt(12)
        dates_para.font.bold = True
        dates_para.font.color.rgb = PRIMARY_COLOR
        dates_para.alignment = PP_ALIGN.CENTER
    
    def add_data_dictionary_slide(self):
        """Add Data Dictionary slide with tables"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Data Dictionary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Create table
        rows, cols = 9, 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(3)
        table.columns[2].width = Inches(3.5)
        
        # Header row
        header_texts = ["Table Name", "Field Name", "Data Type"]
        for col_idx, header_text in enumerate(header_texts):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECONDARY_COLOR
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Data rows
        data = [
            ("Students", "StudentID", "INTEGER (PK)"),
            ("Students", "Name, Email, Class", "VARCHAR, VARCHAR, INT"),
            ("Teachers", "TeacherID", "INTEGER (PK)"),
            ("Teachers", "Name, Subject, Classes", "VARCHAR, VARCHAR, TEXT"),
            ("Attendance", "AttendanceID", "INTEGER (PK)"),
            ("Attendance", "StudentID, Date, Status", "INT (FK), DATE, ENUM"),
            ("Fees", "FeeID", "INTEGER (PK)"),
            ("Fees", "StudentID, Amount, Status", "INT (FK), DECIMAL, ENUM"),
        ]
        
        for row_idx, (table_name, field, data_type) in enumerate(data, 1):
            table.cell(row_idx, 0).text = table_name
            table.cell(row_idx, 1).text = field
            table.cell(row_idx, 2).text = data_type
            
            for col_idx in range(3):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = LIGHT_BG
                else:
                    cell.fill.fore_color.rgb = WHITE
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_COLOR
    
    def add_dfd_slide(self):
        """Add Data Flow Diagram slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Data Flow Diagram (DFD)"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # DFD Entities - simplified visual representation
        # Actor/User
        user_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(2),
            Inches(1.2), Inches(1.2)
        )
        user_shape.fill.solid()
        user_shape.fill.fore_color.rgb = ACCENT_COLOR
        user_shape.line.color.rgb = ACCENT_COLOR
        
        user_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.35),
            Inches(1.2), Inches(0.5)
        )
        user_frame = user_text.text_frame
        user_frame.text = "User/Admin"
        user_para = user_frame.paragraphs[0]
        user_para.font.size = Pt(11)
        user_para.font.bold = True
        user_para.font.color.rgb = WHITE
        user_para.alignment = PP_ALIGN.CENTER
        
        # Process - Application
        process_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3), Inches(1.5),
            Inches(2), Inches(2)
        )
        process_shape.fill.solid()
        process_shape.fill.fore_color.rgb = SECONDARY_COLOR
        process_shape.line.color.rgb = SECONDARY_COLOR
        process_shape.line.width = Pt(2)
        
        process_text = slide.shapes.add_textbox(
            Inches(3.2), Inches(1.8),
            Inches(1.6), Inches(1.4)
        )
        process_frame = process_text.text_frame
        process_frame.word_wrap = True
        process_frame.text = "School Hub\nDashboard\nApplication"
        process_para = process_frame.paragraphs[0]
        process_para.font.size = Pt(12)
        process_para.font.bold = True
        process_para.font.color.rgb = WHITE
        process_para.alignment = PP_ALIGN.CENTER
        
        # Database
        db_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.5), Inches(2),
            Inches(1.2), Inches(1.5)
        )
        db_shape.fill.solid()
        db_shape.fill.fore_color.rgb = RGBColor(156, 39, 176)
        db_shape.line.color.rgb = RGBColor(156, 39, 176)
        
        db_text = slide.shapes.add_textbox(
            Inches(6.4), Inches(2.2),
            Inches(1.4), Inches(1)
        )
        db_frame = db_text.text_frame
        db_frame.word_wrap = True
        db_frame.text = "Supabase\nDatabase"
        db_para = db_frame.paragraphs[0]
        db_para.font.size = Pt(11)
        db_para.font.bold = True
        db_para.font.color.rgb = WHITE
        db_para.alignment = PP_ALIGN.CENTER
        
        # Data flows
        # User to Process
        connector1 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(2.1), Inches(2.4),
            Inches(0.9), Inches(0.4)
        )
        connector1.fill.solid()
        connector1.fill.fore_color.rgb = TEXT_COLOR
        connector1.line.color.rgb = TEXT_COLOR
        
        connector1_text = slide.shapes.add_textbox(
            Inches(2.2), Inches(2.2),
            Inches(0.7), Inches(0.3)
        )
        connector1_frame = connector1_text.text_frame
        connector1_frame.text = "Inputs"
        connector1_para = connector1_frame.paragraphs[0]
        connector1_para.font.size = Pt(9)
        connector1_para.font.color.rgb = TEXT_COLOR
        
        # Process to Database
        connector2 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(5.1), Inches(2.4),
            Inches(1.3), Inches(0.4)
        )
        connector2.fill.solid()
        connector2.fill.fore_color.rgb = TEXT_COLOR
        connector2.line.color.rgb = TEXT_COLOR
        
        connector2_text = slide.shapes.add_textbox(
            Inches(5.3), Inches(2.2),
            Inches(0.9), Inches(0.3)
        )
        connector2_frame = connector2_text.text_frame
        connector2_frame.text = "Queries"
        connector2_para = connector2_frame.paragraphs[0]
        connector2_para.font.size = Pt(9)
        connector2_para.font.color.rgb = TEXT_COLOR
        
        # Database to Process
        connector3 = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW,
            Inches(5.1), Inches(3),
            Inches(1.3), Inches(0.4)
        )
        connector3.fill.solid()
        connector3.fill.fore_color.rgb = ACCENT_COLOR
        connector3.line.color.rgb = ACCENT_COLOR
        
        connector3_text = slide.shapes.add_textbox(
            Inches(5.4), Inches(3.1),
            Inches(0.7), Inches(0.3)
        )
        connector3_frame = connector3_text.text_frame
        connector3_frame.text = "Data"
        connector3_para = connector3_frame.paragraphs[0]
        connector3_para.font.size = Pt(9)
        connector3_para.font.color.rgb = TEXT_COLOR
        
        # Process to User
        connector4 = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_ARROW,
            Inches(2.1), Inches(3),
            Inches(0.9), Inches(0.4)
        )
        connector4.fill.solid()
        connector4.fill.fore_color.rgb = ACCENT_COLOR
        connector4.line.color.rgb = ACCENT_COLOR
        
        connector4_text = slide.shapes.add_textbox(
            Inches(2), Inches(3.1),
            Inches(1), Inches(0.3)
        )
        connector4_frame = connector4_text.text_frame
        connector4_frame.text = "Output"
        connector4_para = connector4_frame.paragraphs[0]
        connector4_para.font.size = Pt(9)
        connector4_para.font.color.rgb = TEXT_COLOR
        
        # Data entities box
        entities_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(9), Inches(2.8)
        )
        entities_frame = entities_box.text_frame
        entities_frame.word_wrap = True
        
        entities_text = """Data Entities:
        • Students: Student profiles, enrollments, contact information
        • Teachers: Teacher profiles, qualifications, subject assignments
        • Classes: Class information, schedules, teacher assignments
        • Attendance: Daily attendance records, patterns
        • Fees: Fee transactions, payment status, balances
        • Users: Admin and staff accounts with role-based access"""
        
        entities_frame.text = entities_text
        for para in entities_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.color.rgb = TEXT_COLOR
            para.space_before = Pt(2)
            para.space_after = Pt(2)
    
    def add_er_diagram_slide(self):
        """Add Entity Relationship Diagram slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Entity-Relationship Diagram (ER Diagram)"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Draw rectangles for entities
        entities = [
            ("Students", Inches(0.5), Inches(1.5), ["StudentID (PK)", "Name", "Email", "Phone", "ClassID (FK)"]),
            ("Teachers", Inches(3.5), Inches(1.5), ["TeacherID (PK)", "Name", "Email", "Subject"]),
            ("Classes", Inches(6.5), Inches(1.5), ["ClassID (PK)", "ClassName", "TeacherID (FK)", "Capacity"]),
            ("Attendance", Inches(0.5), Inches(4), ["AttendanceID (PK)", "StudentID (FK)", "Date", "Status"]),
            ("Fees", Inches(6.5), Inches(4), ["FeeID (PK)", "StudentID (FK)", "Amount", "Status"]),
        ]
        
        for entity_name, x, y, attributes in entities:
            # Entity box
            entity_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y,
                Inches(2.5), Inches(1.8)
            )
            entity_box.fill.solid()
            entity_box.fill.fore_color.rgb = LIGHT_BG
            entity_box.line.color.rgb = SECONDARY_COLOR
            entity_box.line.width = Pt(2)
            
            # Entity name
            entity_text = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.1),
                Inches(2.3), Inches(0.4)
            )
            entity_frame = entity_text.text_frame
            entity_frame.text = entity_name
            entity_para = entity_frame.paragraphs[0]
            entity_para.font.size = Pt(14)
            entity_para.font.bold = True
            entity_para.font.color.rgb = SECONDARY_COLOR
            entity_para.alignment = PP_ALIGN.CENTER
            
            # Attributes
            attr_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.5),
                Inches(2.1), Inches(1.2)
            )
            attr_frame = attr_text.text_frame
            attr_frame.word_wrap = True
            for i, attr in enumerate(attributes):
                if i > 0:
                    attr_frame.add_paragraph()
                p = attr_frame.paragraphs[i]
                p.text = "• " + attr
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(1)
                p.space_after = Pt(1)
        
        # Relationships
        relationships_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2),
            Inches(9), Inches(0.9)
        )
        relationships_frame = relationships_box.text_frame
        relationships_frame.word_wrap = True
        relationships_frame.text = "Relationships: Students → Classes (1:N) | Students → Attendance (1:N) | Students → Fees (1:N) | Teachers → Classes (1:N)"
        relationships_para = relationships_frame.paragraphs[0]
        relationships_para.font.size = Pt(11)
        relationships_para.font.bold = True
        relationships_para.font.color.rgb = PRIMARY_COLOR
        relationships_para.space_before = Pt(5)
    
    def add_interface_slide(self):
        """Add User Interface slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "User Interface Design"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Key Features Grid
        features = [
            ("🎯 Dashboard", "Real-time analytics and metrics"),
            ("👥 Student Mgmt", "Add, edit, track all students"),
            ("🏫 Class Mgmt", "Organize classes & schedules"),
            ("📋 Attendance", "Track daily attendance"),
            ("💰 Fee Tracking", "Monitor fee payments"),
            ("⚙️ Settings", "Configure school parameters"),
        ]
        
        y_pos = 1.5
        for idx, (feature, desc) in enumerate(features):
            col = idx % 2
            row = idx // 2
            
            x_pos = 0.5 + col * 5
            y = y_pos + row * 1.6
            
            # Feature box
            feature_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y),
                Inches(4.5), Inches(1.4)
            )
            feature_box.fill.solid()
            feature_box.fill.fore_color.rgb = LIGHT_BG
            feature_box.line.color.rgb = SECONDARY_COLOR
            feature_box.line.width = Pt(2)
            
            # Feature title
            feature_title = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y + 0.1),
                Inches(4.1), Inches(0.4)
            )
            feature_frame = feature_title.text_frame
            feature_frame.text = feature
            feature_para = feature_frame.paragraphs[0]
            feature_para.font.size = Pt(14)
            feature_para.font.bold = True
            feature_para.font.color.rgb = SECONDARY_COLOR
            
            # Feature description
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y + 0.5),
                Inches(4.1), Inches(0.8)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_frame.text = desc
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(11)
            desc_para.font.color.rgb = TEXT_COLOR
    
    def add_reports_slide(self):
        """Add Expected Reports slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Expected Report Generation"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Reports content
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = LIGHT_BG
        content_box.line.color.rgb = SECONDARY_COLOR
        content_box.line.width = Pt(2)
        
        reports_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(8.4), Inches(5.2)
        )
        reports_frame = reports_text.text_frame
        reports_frame.word_wrap = True
        
        reports_content = """📊 Student Report
        • Individual student profiles with attendance history
        • Performance analytics and fee payment status

        📈 Attendance Report
        • Daily/weekly/monthly attendance statistics
        • Absentee tracking and trends

        💼 Fee Report
        • Outstanding payments by student/class
        • Payment frequency and collection rates

        📚 Class Report
        • Enrollment statistics by class
        • Teacher assignments and class performance metrics

        👨‍🏫 Teacher Report
        • Teacher workload distribution
        • Subject-wise class assignments

        🎯 Dashboard Analytics
        • Real-time overview of all key metrics
        • Graphical visualizations (charts, graphs)
        • Exportable to PDF/CSV formats"""
        
        reports_frame.text = reports_content
        for i, para in enumerate(reports_frame.paragraphs):
            para.font.size = Pt(12)
            para.font.color.rgb = TEXT_COLOR
            if any(char in para.text for char in ['•', '📊', '📈', '💼', '📚', '👨‍🏫', '🎯']):
                para.font.bold = False
                para.space_before = Pt(8)
            else:
                para.space_before = Pt(2)
            para.space_after = Pt(2)
    
    def add_references_slide(self):
        """Add References slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "References & Bibliography"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # References content
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = LIGHT_BG
        content_box.line.color.rgb = SECONDARY_COLOR
        content_box.line.width = Pt(2)
        
        references_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(8.4), Inches(5.2)
        )
        references_frame = references_text.text_frame
        references_frame.word_wrap = True
        
        references_content = """TECHNOLOGIES & FRAMEWORKS:
        1. React 18 - Modern UI library for interactive interfaces
        2. TypeScript - Type-safe JavaScript development
        3. Supabase - Backend-as-a-Service with PostgreSQL
        4. Tailwind CSS - Utility-first CSS framework
        5. shadcn/ui - High-quality React components
        6. Vite - Next-generation frontend build tool
        7. TanStack Query - Server state management
        8. React Router v6 - Client-side routing

        DOCUMENTATION REFERENCES:
        • React Official Documentation: react.dev
        • Supabase Documentation: supabase.com/docs
        • TypeScript Handbook: typescriptlang.org
        • Tailwind CSS Documentation: tailwindcss.com
        • Vite Documentation: vitejs.dev

        DESIGN PATTERNS:
        • MVC Architecture for application structure
        • RESTful API design principles
        • Database normalization (3NF)"""
        
        references_frame.text = references_content
        for para in references_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.color.rgb = TEXT_COLOR
            if any(char in para.text for char in [':', '•']):
                para.font.bold = True
                para.font.color.rgb = SECONDARY_COLOR
            else:
                para.font.bold = False
            para.space_before = Pt(2)
            para.space_after = Pt(2)
    
    def add_future_scope_slide(self):
        """Add Future Scope slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.color.rgb = PRIMARY_COLOR
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Future Scope & Enhancements"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Future features
        features = [
            ("🤖 AI-Powered Insights", ["Predictive analytics for student performance", "Automated attendance anomaly detection"]),
            ("📱 Mobile Application", ["Native iOS & Android apps", "Offline data synchronization"]),
            ("🔐 Advanced Security", ["Two-factor authentication (2FA)", "Role-based access control (RBAC)"]),
            ("📢 Communication Hub", ["In-app messaging system", "Automated SMS/Email notifications"]),
            ("📊 Advanced Analytics", ["Custom report generation", "Data visualization dashboards", "Export to multiple formats"]),
            ("🌐 Multi-Language Support", ["Internationalization (i18n)", "Support for regional languages"]),
        ]
        
        y_pos = 1.5
        for idx, (feature, items) in enumerate(features):
            col = idx % 2
            row = idx // 2
            
            x_pos = 0.5 + col * 5
            y = y_pos + row * 1.95
            
            # Feature box
            feature_height = 0.4 + len(items) * 0.25
            feature_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y),
                Inches(4.5), Inches(feature_height)
            )
            feature_box.fill.solid()
            feature_box.fill.fore_color.rgb = LIGHT_BG
            feature_box.line.color.rgb = SECONDARY_COLOR
            feature_box.line.width = Pt(1.5)
            
            # Feature title
            feature_title = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y + 0.1),
                Inches(4.1), Inches(0.3)
            )
            feature_frame = feature_title.text_frame
            feature_frame.text = feature
            feature_para = feature_frame.paragraphs[0]
            feature_para.font.size = Pt(12)
            feature_para.font.bold = True
            feature_para.font.color.rgb = SECONDARY_COLOR
            
            # Feature items
            items_text = slide.shapes.add_textbox(
                Inches(x_pos + 0.3), Inches(y + 0.4),
                Inches(3.9), Inches(feature_height - 0.5)
            )
            items_frame = items_text.text_frame
            items_frame.word_wrap = True
            for i, item in enumerate(items):
                if i > 0:
                    items_frame.add_paragraph()
                p = items_frame.paragraphs[i]
                p.text = "• " + item
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = Pt(1)
                p.space_after = Pt(1)
    
    def add_conclusion_slide(self):
        """Add Conclusion slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_number += 1
        
        # Background with gradient effect
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Main content area
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1.5),
            Inches(8), Inches(4.5)
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = WHITE
        content_shape.line.color.rgb = SECONDARY_COLOR
        content_shape.line.width = Pt(3)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2),
            Inches(7), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Thank You"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3),
            Inches(7), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "School Hub Dashboard"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.7),
            Inches(7), Inches(1.8)
        )
        contact_frame = contact_box.text_frame
        contact_frame.word_wrap = True
        contact_frame.text = """For Questions & Feedback:
        📧 Project Team
        📱 School Administration
        🔗 GitHub Repository: school-hub-dash-62"""
        
        for para in contact_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = TEXT_COLOR
            para.alignment = PP_ALIGN.CENTER
            para.space_before = Pt(3)
            para.space_after = Pt(3)
    
    def generate(self):
        """Generate the complete presentation"""
        # Slide 1: Title Slide
        self.add_title_slide(
            "School Hub Dashboard",
            "A Comprehensive School Management System",
            "Level-1 PPT Presentation\nFebruary 13-14, 2026"
        )
        
        # Slide 2: Introduction
        self.add_content_slide(
            "Introduction & Project Objectives",
            [
                ("Project Overview", 0),
                ("Modern, full-featured school management system built with React & TypeScript", 1),
                ("Unified platform for managing students, teachers, classes, attendance, and fees", 1),
                ("Key Objectives", 0),
                ("Streamline administrative operations and reduce manual paperwork", 1),
                ("Provide real-time insights into school metrics and student performance", 1),
                ("Enable efficient communication between educators and administrators", 1),
                ("Ensure scalability and security with cloud-based infrastructure", 1),
                ("Target Users", 0),
                ("School administrators, teachers, and office staff", 1),
            ]
        )
        
        # Slide 3: System Requirements
        self.add_two_column_slide(
            "System Requirements (SRS)",
            [
                "✓ Functional Requirements",
                "• Student management (add/edit/delete)",
                "• Attendance tracking system",
                "• Fee collection & tracking",
                "• Teacher & class management",
                "• Dashboard with analytics",
                "• Role-based access control",
                "• Report generation (PDF/CSV)",
            ],
            [
                "✓ Non-Functional Requirements",
                "• Performance: <2s response time",
                "• Security: End-to-end encryption",
                "• Availability: 99.5% uptime",
                "• Scalability: Support 1000+ users",
                "• Compatibility: Chrome, Firefox, Safari",
                "• Data Backup: Daily automated backups",
                "• User Experience: Intuitive UI",
            ]
        )
        
        # Slide 4: Process Logic
        self.add_content_slide(
            "Process Logic & System Flow",
            [
                ("Authentication Process", 0),
                ("1. User enters credentials → System validates with Supabase", 1),
                ("2. Upon success, generates JWT token → User logged in", 1),
                ("Data Management Process", 0),
                ("1. User submits data via form → System validates input", 1),
                ("2. Data sent to Supabase API → Database stores/updates record", 1),
                ("3. System fetches latest data → Displays on dashboard", 1),
                ("Attendance Recording", 0),
                ("1. Mark attendance → Store in database", 1),
                ("2. Generate attendance reports → Visualize trends", 1),
            ]
        )
        
        # Slide 5: Gantt Chart
        self.add_gantt_chart_slide()
        
        # Slide 6: Data Dictionary
        self.add_data_dictionary_slide()
        
        # Slide 7: DFD
        self.add_dfd_slide()
        
        # Slide 8: ER Diagram
        self.add_er_diagram_slide()
        
        # Slide 9: User Interface
        self.add_interface_slide()
        
        # Slide 10: Expected Reports
        self.add_reports_slide()
        
        # Slide 11: References
        self.add_references_slide()
        
        # Slide 12: Future Scope
        self.add_future_scope_slide()
        
        # Slide 13: Conclusion
        self.add_conclusion_slide()
        
        # Save the presentation
        filename = "School_Hub_Dashboard_PPT_Level1.pptx"
        self.prs.save(filename)
        return filename

def main():
    """Main function to generate the PPT"""
    print("🎨 Generating School Hub Dashboard PPT...")
    print("=" * 60)
    
    ppt_generator = SchoolDashboardPPT()
    filename = ppt_generator.generate()
    
    print(f"✅ PPT Generated Successfully!")
    print(f"📁 File: {filename}")
    print(f"📊 Total Slides: {ppt_generator.slide_number}")
    print(f"📅 Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)
    print("\n📝 Notes:")
    print("   • Level-1 Presentation: 13-14 Feb 2026")
    print("   • Level-2 Presentation: 27-28 Feb 2026 (with live demo)")
    print("   • Level-3 Presentation: 11-14 Mar 2026 (with mentor feedback)")
    print("   • Level-4 Submission: 25-28 Mar 2026 (final report)")
    print("\n✨ Presentation is ready for your Level-1 PPT presentation!")

if __name__ == "__main__":
    main()
